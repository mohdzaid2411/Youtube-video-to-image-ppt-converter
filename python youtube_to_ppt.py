import os
import subprocess
import sys
import cv2
from pptx import Presentation
from pptx.util import Inches

# Function to download YouTube video using yt-dlp
def download_youtube_video(video_url, output_path):
    # Use yt-dlp directly via Python without format selection
    command = [sys.executable, '-m', 'yt_dlp', '-o', f'{output_path}/%(title)s.%(ext)s', video_url]
    subprocess.run(command, check=True)
    
    # Find the downloaded video file (assuming mp4)
    for file in os.listdir(output_path):
        if file.endswith(".mp4"):
            return os.path.join(output_path, file)
    return None

# Function to extract frames from the video
def extract_frames(video_file, output_dir, frame_interval=5):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    cap = cv2.VideoCapture(video_file)
    fps = int(cap.get(cv2.CAP_PROP_FPS))
    count = 0
    extracted_frames = []
    
    while True:
        ret, frame = cap.read()
        if not ret:
            break
        
        if count % (frame_interval * fps) == 0:
            frame_filename = os.path.join(output_dir, f"frame_{count}.jpg")
            cv2.imwrite(frame_filename, frame)
            extracted_frames.append(frame_filename)
        
        count += 1
    
    cap.release()
    return extracted_frames

# Function to create PowerPoint from images
def create_ppt_from_images(image_files, output_ppt):
    prs = Presentation()
    
    for image_file in image_files:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(image_file, Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))
    
    prs.save(output_ppt)

# Main function
def main():
    video_url = input("Enter YouTube video URL: ")
    output_dir = "extracted_frames"
    
    try:
        # Step 1: Download YouTube video using yt-dlp
        print("Downloading video...")
        video_file = download_youtube_video(video_url, output_dir)
        print(f"Video downloaded successfully! {video_file}")
        
        # Step 2: Extract frames
        print("Extracting frames from video...")
        frames = extract_frames(video_file, output_dir)
        print(f"Extracted {len(frames)} frames.")
        
        # Step 3: Create PPT from frames
        print("Creating PowerPoint presentation...")
        create_ppt_from_images(frames, "youtube_video_slides.pptx")
        print("PowerPoint created successfully!")
    
    except Exception as e:
        print(f"Error: {str(e)}")

if __name__ == "__main__":
    main()
